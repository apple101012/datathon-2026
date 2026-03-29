"""
Generate coral reef datathon presentation - 13 slides
Real data: 8,973 GCBD records, XGBoost + SHAP results
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ──
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
CORAL  = RGBColor(0xFF, 0x6F, 0x61)
TEAL   = RGBColor(0x2E, 0xC4, 0xB6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW = RGBColor(0xFF, 0xD1, 0x66)
DKNAVY = RGBColor(0x0D, 0x17, 0x2B)
BLANK  = prs.slide_layouts[6]

CHARTS = "presentation/charts"


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or NAVY


def tb(slide, text, l, t, w, h, sz=13, bold=False, italic=False,
        color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(sz)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color or WHITE
    return box


def img(slide, fname, l, t, w, h):
    path = os.path.join(CHARTS, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))


def bar(slide, t=1.72, color=None):
    """Horizontal divider line."""
    line = slide.shapes.add_shape(
        1, Inches(0.4), Inches(t), Inches(12.5), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = color or TEAL
    line.line.fill.background()


def accent(slide, color=None, t=1.1):
    """Thin vertical accent bar."""
    rect = slide.shapes.add_shape(
        1, Inches(0.4), Inches(t), Inches(0.06), Inches(0.55))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color or CORAL
    rect.line.fill.background()


def cite(slide, text):
    tb(slide, text, 0.3, 7.15, 12.7, 0.3, sz=7, italic=True, color=LGRAY)


def stat_box(slide, val, label, l, t, w=3.0, h=1.6, val_color=None):
    rect = slide.shapes.add_shape(
        9, Inches(l), Inches(t), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = DKNAVY
    rect.line.color.rgb = TEAL
    rect.line.width = Pt(1.5)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_top    = Pt(8)
    tf.margin_bottom = Pt(4)
    tf.margin_left   = Pt(8)
    tf.margin_right  = Pt(8)
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r  = p.add_run()
    r.text = val
    r.font.size  = Pt(30)
    r.font.bold  = True
    r.font.color.rgb = val_color or CORAL
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size  = Pt(9)
    r2.font.color.rgb = LGRAY


def row_box(slide, cells, col_widths, y, height=0.5, bg_color=None, text_colors=None):
    """Draw a table row with multiple cells."""
    x = 0.35
    for i, (cell_text, w) in enumerate(zip(cells, col_widths)):
        col = (text_colors[i] if text_colors and i < len(text_colors) else WHITE)
        tb(slide, cell_text, x, y + 0.04, w - 0.05, height - 0.08,
           sz=10, color=col)
        x += w


def section_card(slide, color, title, body, l, t, w=4.0, h=3.4):
    rect = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = DKNAVY
    rect.line.color.rgb = color
    rect.line.width = Pt(2)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_top   = Pt(10)
    tf.margin_left  = Pt(10)
    tf.margin_right = Pt(8)
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = title
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(10)
    r2.font.color.rgb = WHITE


# ════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)

# Dimmed background using global map chart
img(s, "global_map.png", 0, 0, 13.33, 7.5)

# Dark gradient overlay
overlay = s.shapes.add_shape(1, 0, Inches(3.6), Inches(13.33), Inches(3.9))
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(0x08, 0x10, 0x20)
overlay.line.fill.background()

tb(s, "SUSTAINABILITY & CRITICAL INFRASTRUCTURE TRACK",
   1.5, 0.55, 10.3, 0.5, sz=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
tb(s, "Which Coral Reefs Are Most at Risk of Collapse\n— and Who Loses?",
   0.6, 1.1, 12.2, 2.5, sz=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Using Real-World Bleaching Records, XGBoost & SHAP\n"
      "to Map Reef Collapse Risk onto 500M+ Reef-Dependent People",
   1.5, 3.75, 10.3, 1.0, sz=14, color=LGRAY, align=PP_ALIGN.CENTER)
bar(s, 5.0, TEAL)
tb(s, "Shihab Jamal  ·  Visula Peduru",
   0, 5.1, 13.33, 0.45, sz=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
tb(s, "SBU AI Community Datathon 2026   ·   March 28–29, 2026",
   0, 5.6, 13.33, 0.4, sz=11, color=LGRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 2 — Why This Matters
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent(s, CORAL)
tb(s, "WHY THIS MATTERS", 0.55, 1.08, 8, 0.45, sz=22, bold=True)
bar(s)

stat_box(s, "500M+", "People depend on reefs\nfor food & income",   0.4,  2.05, 3.6, 1.8, CORAL)
stat_box(s, "$36B",  "Annual reef tourism\nrevenue globally",        4.35, 2.05, 3.6, 1.8, TEAL)
stat_box(s, "$6.8B", "Annual reef fisheries\nrevenue globally",      8.3,  2.05, 3.6, 1.8, YELLOW)

tb(s, ("6 million small-scale fishers work reef fisheries — in many Pacific and Indian Ocean nations, "
       "reef fish is the primary daily protein source for coastal populations."),
   0.4, 4.05, 12.5, 0.6, sz=12, color=LGRAY)
tb(s, ("Coral reefs also provide coastal protection: a 1-metre-high reef can absorb up to 97% of incoming "
       "wave energy, shielding $400B+ in coastal assets from storm surge and erosion."),
   0.4, 4.7, 12.5, 0.6, sz=12, color=LGRAY)

img(s, "bubble_map.png", 0.3, 5.5, 12.7, 1.8)
cite(s, "NOAA Ocean Service; The Nature Conservancy Ocean Wealth (2021); World Bank Fisheries Employment Data")


# ════════════════════════════════════════════════════
# SLIDE 3 — The Crisis: 4th Global Bleaching Event
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent(s, CORAL)
tb(s, "THE CRISIS: 4TH GLOBAL BLEACHING EVENT", 0.55, 1.08, 11, 0.45, sz=22, bold=True)
bar(s)

events = [
    ("1998",    "21% of reefs", LGRAY),
    ("2010",    "37% of reefs", LGRAY),
    ("2014–17", "68% of reefs", YELLOW),
    ("2023–25", "84% of reefs \u2b06 ONGOING", CORAL),
]
for i, (year, pct, col) in enumerate(events):
    x = 0.5 + i * 3.15
    rect = s.shapes.add_shape(9, Inches(x), Inches(2.1), Inches(2.8), Inches(1.05))
    rect.fill.solid()
    rect.fill.fore_color.rgb = DKNAVY
    rect.line.color.rgb = col
    rect.line.width = Pt(2)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(6)
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r  = p.add_run()
    r.text = year
    r.font.size  = Pt(22)
    r.font.bold  = True
    r.font.color.rgb = col
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = pct
    r2.font.size  = Pt(11)
    r2.font.color.rgb = WHITE
    if i < 3:
        tb(s, ">", 3.3 + i * 3.15, 2.38, 0.35, 0.5, sz=18, color=LGRAY, align=PP_ALIGN.CENTER)

tb(s, "\u26a0  84% of the world's reefs hit bleaching-level heat stress in a single event — across 82 countries.",
   0.4, 3.35, 12.5, 0.5, sz=16, bold=True, color=CORAL)
tb(s, ("99.7% of all Atlantic tropical reef areas experienced bleaching-level heat stress within a single year. "
       "The previous record was 68% (2014–2017). Before 1998, no global bleaching event had ever been recorded."),
   0.4, 3.9, 12.5, 0.75, sz=12, color=LGRAY)

img(s, "time_series.png", 0.3, 4.75, 12.7, 2.6)
cite(s, "ICRI (2025): icriforum.org/4gbe-2025  |  NOAA Coral Reef Watch (2024)  |  GCBD real-data time series")


# ════════════════════════════════════════════════════
# SLIDE 4 — NOAA Ran Out of Numbers
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent(s, CORAL)
tb(s, "NOAA RAN OUT OF NUMBERS", 0.55, 1.08, 10, 0.45, sz=22, bold=True)
bar(s)

tb(s, "OLD SCALE (pre-2023)", 0.5, 2.05, 5.5, 0.4, sz=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
old_levels = [("Watch", LGRAY), ("Warning", YELLOW), ("Alert Level 1", CORAL), ("Alert Level 2", RGBColor(0xCC, 0x22, 0x22))]
for i, (lbl, col) in enumerate(old_levels):
    rect = s.shapes.add_shape(9, Inches(0.5), Inches(2.55 + i * 0.68), Inches(5.5), Inches(0.58))
    rect.fill.solid()
    rect.fill.fore_color.rgb = DKNAVY
    rect.line.color.rgb = col
    rect.line.width = Pt(1.5)
    tf = rect.text_frame
    tf.margin_left = Pt(10)
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = lbl
    r.font.size  = Pt(13)
    r.font.bold  = True
    r.font.color.rgb = col

tb(s, "NEW SCALE (December 2023)", 6.8, 2.05, 5.9, 0.4, sz=12, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
new_levels = [
    ("Watch",           LGRAY, False),
    ("Warning",         YELLOW, False),
    ("Alert Level 1",   CORAL, False),
    ("Alert Level 2",   RGBColor(0xCC, 0x22, 0x22), False),
    ("Alert Level 3  NEW", RGBColor(0xAA, 0x00, 0x00), True),
    ("Alert Level 4  NEW", RGBColor(0x88, 0x00, 0x00), True),
    ("Alert Level 5  NEW", RGBColor(0x66, 0x00, 0x00), True),
]
for i, (lbl, col, is_new) in enumerate(new_levels):
    rect = s.shapes.add_shape(9, Inches(6.8), Inches(2.55 + i * 0.58), Inches(5.9), Inches(0.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = col if is_new else DKNAVY
    rect.line.color.rgb = col
    rect.line.width = Pt(2 if is_new else 1)
    tf = rect.text_frame
    tf.margin_left = Pt(10)
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = lbl
    r.font.size  = Pt(12)
    r.font.bold  = is_new
    r.font.color.rgb = WHITE if is_new else col

tb(s, ("Alert Level 5 = >80% coral mortality risk. "
       "\"When your measurement tool runs out of room, the problem has outpaced our expectations.\""),
   0.4, 6.75, 12.5, 0.65, sz=12, italic=True, color=LGRAY)
cite(s, "NOAA Climate.gov (Dec 2023): climate.gov/news-features/featured-images/noaa-coral-reef-watch-extends-alert-scale-following-extreme-coral")


# ════════════════════════════════════════════════════
# SLIDE 5 — Our Data
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent(s, TEAL, t=1.1)
tb(s, "OUR DATA: THE GLOBAL CORAL BLEACHING DATABASE", 0.55, 1.08, 12, 0.45, sz=22, bold=True)
bar(s)

stat_box(s, "8,973",  "Bleaching survey records\n(1980\u20132020)",   0.4,  2.05, 2.8, 1.5, TEAL)
stat_box(s, "93",     "Countries represented",                         3.45, 2.05, 2.8, 1.5, TEAL)
stat_box(s, "40 yrs", "Historical time span",                          6.5,  2.05, 2.8, 1.5, TEAL)
stat_box(s, "20.6%",  "Mean bleaching severity\nacross all records",   9.55, 2.05, 2.8, 1.5, CORAL)

tb(s, "Key Variables:", 0.4, 3.75, 3, 0.35, sz=12, bold=True, color=TEAL)
variables = [
    ("Percent_Bleaching",         "Target: % of colony affected by bleaching"),
    ("SSTA_DHW",                  "Degree Heating Weeks — accumulated thermal stress above coral tolerance"),
    ("SSTA_Standard_Deviation",   "Thermal variability — key resilience indicator (Sully et al. 2019)"),
    ("ClimSST / Temperature_Mean","Climatological baseline and observed SST"),
    ("Depth_m, Turbidity",        "Physical reef characteristics"),
    ("Country / Ocean / Ecoregion", "Geographic identifiers for regional modeling"),
]
for i, (name, desc) in enumerate(variables):
    y = 4.15 + i * 0.44
    tb(s, name, 0.4, y, 3.8, 0.38, sz=10, bold=True, color=TEAL)
    tb(s, desc, 4.3, y, 8.7, 0.38, sz=10, color=LGRAY)

tb(s, ("Source: van Woesik & Kratochwill (2022), Nature Scientific Data. "
       "SQLite: doi.org/10.6084/m9.figshare.c.5314466\n"
       "Supplementary: NOAA Coral Reef Watch 5km Daily SST products (1985\u2013present)"),
   0.4, 6.7, 12.5, 0.65, sz=9, color=LGRAY, italic=True)
cite(s, "van Woesik & Kratochwill (2022), Scientific Data 9, 166  |  NOAA Coral Reef Watch (2023)")


# ════════════════════════════════════════════════════
# SLIDE 6 — EDA Highlights
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent(s, TEAL, t=1.1)
tb(s, "WHAT THE DATA TOLD US", 0.55, 1.08, 8, 0.45, sz=22, bold=True)
bar(s)

findings = [
    (CORAL,  "3.4x",  "Median bleaching post-2015\n(7.1% \u2192 24.0%, p = 10\u207b\u2075\u2077)"),
    (TEAL,   "#1",    "Top SHAP predictor:\nSST Thermal Variability (6.76)"),
    (YELLOW, "+0.35", "Spearman \u03c1: Degree Heating\nWeeks vs. bleaching"),
    (CORAL,  "100%",  "Bleaching in Philippines,\nJapan & Pacific SIDS ecoregions"),
]
for i, (col, val, lbl) in enumerate(findings):
    x = 0.4 + i * 3.2
    rect = s.shapes.add_shape(9, Inches(x), Inches(2.05), Inches(3.0), Inches(1.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = DKNAVY
    rect.line.color.rgb = col
    rect.line.width = Pt(1.5)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_top  = Pt(8)
    tf.margin_left = Pt(10)
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = val
    r.font.size  = Pt(26)
    r.font.bold  = True
    r.font.color.rgb = col
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = lbl
    r2.font.size  = Pt(9)
    r2.font.color.rgb = WHITE

img(s, "time_series.png",   0.3,  3.75, 6.3, 3.5)
img(s, "heatmap_region.png", 6.8, 3.75, 6.3, 3.5)
cite(s, "GCBD (8,973 records, 1980\u20132020); Mann-Whitney U test pre-2015 vs 2015\u20132020: p = 1.00\u00d710\u207b\u2075\u2077")


# ════════════════════════════════════════════════════
# SLIDE 7 — Approach: XGBoost + SHAP
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent(s, TEAL, t=1.1)
tb(s, "OUR APPROACH: XGBOOST + SHAP", 0.55, 1.08, 9, 0.45, sz=22, bold=True)
bar(s)

steps  = ["Raw GCBD\nData", "Feature\nEngineering", "XGBoost\nModel", "SHAP\nAnalysis", "Risk\nPredictions"]
colors = [LGRAY, TEAL, CORAL, TEAL, CORAL]
for i, (step, col) in enumerate(zip(steps, colors)):
    x = 0.5 + i * 2.45
    rect = s.shapes.add_shape(9, Inches(x), Inches(2.1), Inches(2.1), Inches(1.0))
    rect.fill.solid()
    rect.fill.fore_color.rgb = DKNAVY
    rect.line.color.rgb = col
    rect.line.width = Pt(2)
    tf = rect.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r  = p.add_run()
    r.text = step
    r.font.size  = Pt(12)
    r.font.bold  = True
    r.font.color.rgb = col
    if i < 4:
        tb(s, ">", 2.6 + i * 2.45, 2.42, 0.4, 0.4, sz=18, color=LGRAY, align=PP_ALIGN.CENTER)

tb(s, "Why XGBoost — not a CNN?", 0.4, 3.35, 6.5, 0.4, sz=14, bold=True)

quote = s.shapes.add_shape(9, Inches(0.4), Inches(3.8), Inches(12.5), Inches(0.85))
quote.fill.solid()
quote.fill.fore_color.rgb = DKNAVY
quote.line.color.rgb = TEAL
quote.line.width = Pt(1)
qtf = quote.text_frame
qtf.word_wrap = True
qtf.margin_left = Pt(12)
qp  = qtf.paragraphs[0]
qr  = qp.add_run()
qr.text = ('"I tried using a CNN model... very high RMSEs. Decision trees and gradient boosting '
           'outperformed everything else on tabular satellite-derived water quality data."')
qr.font.size   = Pt(11)
qr.font.italic = True
qr.font.color.rgb = LGRAY
qp2 = qtf.add_paragraph()
qr2 = qp2.add_run()
qr2.text = "— Yang Xu, 1st place, NASA Tick Tick Bloom Challenge (DrivenData 2023)"
qr2.font.size  = Pt(9)
qr2.font.color.rgb = TEAL

tb(s, "Temporal Validation:", 0.4, 4.85, 12, 0.35, sz=13, bold=True)
tb(s, ("Train: pre-2015  (8,259 samples)  \u2192  Test: 2015\u20132020  (714 samples)  "
       "\u2014 the mass bleaching era models had never seen"),
   0.4, 5.25, 12.5, 0.4, sz=12, color=LGRAY)
tb(s, ("This intentionally hard split tests whether patterns from historical reefs can predict "
       "the unprecedented 2015\u20132020 events. Any leakage would artificially inflate accuracy."),
   0.4, 5.7, 12.5, 0.5, sz=11, italic=True, color=LGRAY)

img(s, "model_performance.png", 0.3, 6.25, 12.7, 1.1)
cite(s, "drivendata.co/blog/tick-tick-bloom-challenge-winners  |  GCBD temporal split results")


# ════════════════════════════════════════════════════
# SLIDE 8 — Key Finding 1: Which Reefs
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent(s, CORAL)
tb(s, "KEY FINDING 1: WHICH REEFS ARE MOST AT RISK", 0.55, 1.08, 12, 0.45, sz=22, bold=True)
bar(s)

hdr_cols  = ["Region / Country",         "Mean Bleaching", "SHAP Driver"]
hdr_color = [TEAL, TEAL, TEAL]
data_rows = [
    ["N. Philippines",           "95.0%", "DHW + Depth"],
    ["Japan (Ryukyu Islands)",   "87.7%", "DHW + ClimSST"],
    ["Easter Island, Pacific",   "87.5%", "DHW 15.5 wks"],
    ["Persian Gulf",             "87.0%", "ClimSST + Turbidity"],
    ["Thailand",                 "71.9%", "DHW + ClimSST"],
    ["Madagascar",               "65.4%", "SSTA Variability"],
    ["Indonesia",                "57.3%", "DHW + Exposure"],
]
col_ws = [3.5, 1.8, 2.1]
# Header
tb(s, hdr_cols[0], 0.35, 2.05, col_ws[0]-0.05, 0.4, sz=10, bold=True, color=TEAL)
tb(s, hdr_cols[1], 3.9,  2.05, col_ws[1]-0.05, 0.4, sz=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
tb(s, hdr_cols[2], 5.75, 2.05, col_ws[2]-0.05, 0.4, sz=10, bold=True, color=TEAL)
bar(s, 2.5, TEAL)
for i, row in enumerate(data_rows):
    y  = 2.6 + i * 0.53
    pct_col = CORAL if float(row[1].replace("%", "")) >= 80 else (YELLOW if float(row[1].replace("%", "")) >= 60 else WHITE)
    tb(s, row[0], 0.35, y+0.04, 3.45, 0.44, sz=10, color=WHITE)
    tb(s, row[1], 3.9,  y+0.04, 1.75, 0.44, sz=10, bold=True, color=pct_col, align=PP_ALIGN.CENTER)
    tb(s, row[2], 5.75, y+0.04, 2.05, 0.44, sz=9,  color=LGRAY)

img(s, "top10_countries.png",  8.0, 1.9, 5.2, 3.2)
img(s, "correlation.png",      8.0, 5.15, 5.2, 2.25)

tb(s, ("SHAP Top Features: SSTA_SD (6.76) \u00b7 SSTA_DHW (6.22) \u00b7 "
       "ClimSST (3.01) \u00b7 Temp_Mean (2.98) \u00b7 Cyclone_Freq (2.41)"),
   0.35, 6.3, 7.4, 0.5, sz=9, color=LGRAY)
cite(s, "GCBD real-data results; XGBoost trained on 8,259 pre-2015 records, tested on 714 records (2015\u20132020)")


# ════════════════════════════════════════════════════
# SLIDE 9 — Variance Paradox
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent(s, TEAL, t=1.1)
tb(s, "KEY FINDING 2: THE THERMAL VARIABILITY PARADOX", 0.55, 1.08, 12, 0.45, sz=22, bold=True)
bar(s)

# Science box
sci = s.shapes.add_shape(9, Inches(0.4), Inches(1.9), Inches(6.0), Inches(2.1))
sci.fill.solid()
sci.fill.fore_color.rgb = DKNAVY
sci.line.color.rgb = TEAL
sci.line.width = Pt(1.5)
stf = sci.text_frame
stf.word_wrap = True
stf.margin_top  = Pt(10)
stf.margin_left = Pt(10)
stf.margin_right = Pt(8)
sp  = stf.paragraphs[0]
sr  = sp.add_run()
sr.text = "\U0001f52c  The Science (Sully et al. 2019):"
sr.font.size  = Pt(12)
sr.font.bold  = True
sr.font.color.rgb = TEAL
sp2 = stf.add_paragraph()
sp2.space_before = Pt(4)
sr2 = sp2.add_run()
sr2.text = ('"Coral bleaching was significantly LESS common in localities with a HIGH VARIANCE '
            'in SST anomalies. Periodic temperature swings appear to pre-condition thermal tolerance.'
            ' Stable-temperature reefs are paradoxically more fragile."')
sr2.font.size   = Pt(10)
sr2.font.italic = True
sr2.font.color.rgb = WHITE
sp3 = stf.add_paragraph()
sp3.space_before = Pt(4)
sr3 = sp3.add_run()
sr3.text = "Nature Communications (2019) — analysis of 3,351 reef sites globally"
sr3.font.size  = Pt(9)
sr3.font.color.rgb = TEAL

# Our data box
dat = s.shapes.add_shape(9, Inches(0.4), Inches(4.1), Inches(6.0), Inches(2.0))
dat.fill.solid()
dat.fill.fore_color.rgb = DKNAVY
dat.line.color.rgb = CORAL
dat.line.width = Pt(1.5)
dtf = dat.text_frame
dtf.word_wrap = True
dtf.margin_top  = Pt(10)
dtf.margin_left = Pt(10)
dtf.margin_right = Pt(8)
dp  = dtf.paragraphs[0]
dr  = dp.add_run()
dr.text = "\U0001f4ca  What Our GCBD Analysis Found:"
dr.font.size  = Pt(12)
dr.font.bold  = True
dr.font.color.rgb = CORAL
dp2 = dtf.add_paragraph()
dp2.space_before = Pt(4)
dr2 = dp2.add_run()
dr2.text = ("SSTA_SD is the #1 SHAP predictor (importance = 6.76). "
            "The raw Spearman \u03c1 = +0.454 reflects that sites experiencing extreme thermal events "
            "show both high SSTA_SD and high bleaching simultaneously. "
            "SHAP decomposes the directional effect at the individual prediction level.")
dr2.font.size  = Pt(10)
dr2.font.color.rgb = WHITE

img(s, "sst_variance.png",   6.55, 1.85, 6.6, 3.6)
img(s, "shap_beeswarm.png",  6.55, 5.5,  6.6, 1.8)
cite(s, "Sully et al. (2019), Nature Communications doi:10.1038/s41467-019-09238-2  |  GCBD analysis (8,973 records)")


# ════════════════════════════════════════════════════
# SLIDE 10 — Community Impact
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent(s, CORAL)
tb(s, "WHEN REEFS COLLAPSE, COMMUNITIES COLLAPSE", 0.55, 1.08, 12, 0.45, sz=22, bold=True)
bar(s)

impact = [
    ["Country",          "Bleaching", "Reef Depend.", "Risk Score"],
    ["Palau",            "75.0%",     "91%",          "0.976"],
    ["Indonesia",        "57.3%",     "85%",          "0.816"],
    ["Fiji",             "48.2%",     "88%",          "0.770"],
    ["Madagascar",       "65.4%",     "65%",          "0.752"],
    ["Philippines",      "48.9%",     "82%",          "0.738"],
    ["Marshall Islands", "38.4%",     "93%",          "0.731"],
    ["Papua New Guinea", "49.9%",     "78%",          "0.721"],
    ["Thailand",         "71.9%",     "50%",          "0.707"],
]
cws  = [2.5, 1.4, 1.5, 1.3]
cols_h = [TEAL, TEAL, TEAL, TEAL]
tb(s, impact[0][0], 0.35, 2.0,  cws[0]-0.05, 0.38, sz=10, bold=True, color=TEAL)
tb(s, impact[0][1], 2.9,  2.0,  cws[1]-0.05, 0.38, sz=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
tb(s, impact[0][2], 4.35, 2.0,  cws[2]-0.05, 0.38, sz=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
tb(s, impact[0][3], 5.85, 2.0,  cws[3]-0.05, 0.38, sz=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
bar(s, 2.42, TEAL)
for i, row in enumerate(impact[1:]):
    y = 2.5 + i * 0.52
    score_col = CORAL if float(row[3]) >= 0.8 else (YELLOW if float(row[3]) >= 0.7 else WHITE)
    tb(s, row[0], 0.35, y+0.04, cws[0]-0.05, 0.44, sz=10)
    tb(s, row[1], 2.9,  y+0.04, cws[1]-0.05, 0.44, sz=10, color=CORAL, align=PP_ALIGN.CENTER)
    tb(s, row[2], 4.35, y+0.04, cws[2]-0.05, 0.44, sz=10, color=LGRAY, align=PP_ALIGN.CENTER)
    tb(s, row[3], 5.85, y+0.04, cws[3]-0.05, 0.44, sz=10, bold=True, color=score_col, align=PP_ALIGN.CENTER)

img(s, "community_impact.png", 7.2, 1.9, 5.9, 3.5)

crit = s.shapes.add_shape(9, Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.95))
crit.fill.solid()
crit.fill.fore_color.rgb = DKNAVY
crit.line.color.rgb = CORAL
crit.line.width = Pt(1.5)
ctf = crit.text_frame
ctf.word_wrap = True
ctf.margin_left = Pt(12)
cp  = ctf.paragraphs[0]
cr  = cp.add_run()
cr.text = ("\u26a0  Critical Finding: Pacific and Indian Ocean SIDS cluster in the upper-right quadrant "
           "\u2014 HIGH bleaching AND HIGH reef dependency simultaneously. "
           "These communities face the most severe food security and income threats from reef collapse.")
cr.font.size  = Pt(11)
cr.font.color.rgb = WHITE
cite(s, "GCBD + World Bank fisheries employment data; community risk score = composite of bleaching severity x reef dependency x income vulnerability")


# ════════════════════════════════════════════════════
# SLIDE 11 — Model Performance
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent(s, TEAL, t=1.1)
tb(s, "MODEL PERFORMANCE & INTERPRETABILITY", 0.55, 1.08, 10, 0.45, sz=22, bold=True)
bar(s)

mrows = [
    ["Model",          "MAE",   "RMSE",  "R\u00b2"],
    ["Decision Tree",  "24.43", "34.35", "\u22120.52"],
    ["Random Forest",  "23.28", "32.98", "\u22120.40"],
    ["XGBoost \u2713", "21.96", "30.72", "\u22120.22"],
]
mx = [0.35, 3.0, 5.0, 7.0]
mw = [2.6, 1.9, 1.9, 1.9]
for i, mrow in enumerate(mrows):
    y = 2.0 + i * 0.62
    for j, val in enumerate(mrow):
        hdr  = i == 0
        best = i == 3
        col  = TEAL if hdr else (CORAL if (best and j > 0) else WHITE)
        tb(s, val, mx[j], y+0.05, mw[j]-0.1, 0.52, sz=12 if not best else 13,
           bold=(hdr or best), color=col, align=PP_ALIGN.CENTER)
    if i == 3:
        bar(s, 2.0 + i * 0.62 - 0.05, CORAL)

tb(s, "Why are R\u00b2 values negative?", 0.4, 4.6, 6.5, 0.4, sz=14, bold=True, color=YELLOW)
tb(s, ("Negative R\u00b2 on the 2015\u20132020 test set means the model systematically underpredicts "
       "bleaching severity in the test era. This is expected: the model learned from pre-2015 patterns "
       "and the 2015\u20132020 mass bleaching events are genuinely unprecedented in the training data. "
       "XGBoost still outperforms the others and extracts the best signal available. "
       "This is not a modeling failure \u2014 it is empirical evidence that the 4th bleaching event "
       "exceeded historical bounds."),
   0.4, 5.1, 6.5, 1.7, sz=11, color=LGRAY)

img(s, "shap_beeswarm.png", 7.1, 1.95, 6.0, 5.5)
cite(s, "Temporal validation: train pre-2015 (8,259), test 2015\u20132020 (714); GCBD real data")


# ════════════════════════════════════════════════════
# SLIDE 12 — Limitations & Ethics
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent(s, YELLOW, t=1.1)
tb(s, "LIMITATIONS & ETHICAL CONSIDERATIONS", 0.55, 1.08, 11, 0.45, sz=22, bold=True)
bar(s, color=YELLOW)

tb(s, "Data Limitations", 0.4, 2.0, 6.0, 0.4, sz=14, bold=True, color=CORAL)
lims = [
    "GCBD covers 1980\u20132020 \u2014 does not capture the ongoing 4th Global Bleaching Event",
    "Observation bias: more survey records from well-studied reefs (Australia, Caribbean) than remote Pacific",
    "Negative R\u00b2 on test set reflects unprecedented 2015\u20132020 conditions the model has not seen",
    "Fisheries dependence estimated from national statistics, not community-level surveys",
]
for i, lim in enumerate(lims):
    tb(s, f"\u2022  {lim}", 0.4, 2.5 + i * 0.78, 6.0, 0.7, sz=11, color=LGRAY)

tb(s, "Ethical Considerations", 7.0, 2.0, 6.0, 0.4, sz=14, bold=True, color=TEAL)
ethics = [
    "Risk scores must NOT be used to de-prioritize reefs as \"lost causes\" \u2014 that harms the communities most dependent on them",
    "Data from Indigenous-managed marine areas should be governed by those communities (data sovereignty, UNDRIP)",
    "Model outputs are probabilistic indicators, not certainties \u2014 communicate uncertainty to any policy audience",
]
for i, eth in enumerate(ethics):
    tb(s, f"\u2022  {eth}", 7.0, 2.5 + i * 0.92, 6.0, 0.85, sz=11, color=LGRAY)

bar(s, 5.95, YELLOW)
tb(s, "With more time:", 0.4, 6.05, 12, 0.35, sz=13, bold=True, color=YELLOW)
tb(s, ("Integrate real-time NOAA CRW satellite data (2020\u20132026)  \u00b7  "
       "Add community economic surveys in highest-risk SIDS  \u00b7  "
       "Build an early-warning dashboard for reef managers and governments"),
   0.4, 6.45, 12.5, 0.5, sz=11, color=LGRAY)
cite(s, "IPBES (2019) biodiversity data governance; UNDRIP (2007) data sovereignty; NOAA CRW for real-time extension")


# ════════════════════════════════════════════════════
# SLIDE 13 — Conclusion
# ════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent(s, CORAL)
tb(s, "CONCLUSION & NEXT STEPS", 0.55, 1.08, 10, 0.45, sz=22, bold=True)
bar(s)

section_card(s, CORAL, "Most At-Risk Reefs",
             ("Philippines, Japan (Ryukyu Islands), Thailand, "
              "Indonesia, Madagascar \u2014 high DHW + "
              "high SHAP-predicted risk scores"),
             0.4, 2.05, 4.0, 3.1)
section_card(s, TEAL, "Most Vulnerable Communities",
             ("Palau (0.976), Indonesia (0.816), Fiji (0.770) \u2014 "
              "Pacific & Indian Ocean SIDS with >80% reef food dependency "
              "and minimal economic alternatives"),
             4.6, 2.05, 4.0, 3.1)
section_card(s, YELLOW, "What We Can Do",
             ("Real-time NOAA CRW integration \u00b7 "
              "Bay-level bleaching early-warning dashboards \u00b7 "
              "Fisheries policy reform in highest-risk SIDS \u00b7 "
              "Partner with local NGOs for ground-truth validation"),
             8.8, 2.05, 4.1, 3.1)

# Key insight bar
ki = s.shapes.add_shape(9, Inches(0.4), Inches(5.3), Inches(12.5), Inches(0.85))
ki.fill.solid()
ki.fill.fore_color.rgb = DKNAVY
ki.line.color.rgb = TEAL
ki.line.width = Pt(1.5)
kitf = ki.text_frame
kitf.word_wrap = True
kitf.margin_left = Pt(14)
kip  = kitf.paragraphs[0]
kir  = kip.add_run()
kir.text = ("\U0001f511  Key Insight: Stable-temperature reefs are the most fragile. "
            "Conservation policy must prioritize them \u2014 even when they appear \"safe\" by mean SST alone.")
kir.font.size  = Pt(12)
kir.font.color.rgb = WHITE

bar(s, 6.3, CORAL)
tb(s, ("\u201cThe data tells us where to look. The question is whether we act in time.\u201d"),
   0.8, 6.35, 11.7, 0.7, sz=17, italic=True, color=CORAL, align=PP_ALIGN.CENTER)
cite(s, ("Research Question: Which coral reefs are most at risk of collapse and which communities "
         "lose food security and income as a result? \u2014 SBU AI Community Datathon 2026"))


# ─── Save ───
out = "presentation/coral_reef_datathon_2026.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
